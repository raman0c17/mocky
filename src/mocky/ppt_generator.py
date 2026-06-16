"""Cross-platform PowerPoint generation backed by python-pptx.

This engine produces real ``.pptx`` files on Windows (x64/ARM), macOS
(Intel/Apple Silicon) and Linux without requiring Microsoft PowerPoint to be
installed. It is the default Mocky backend and replaces the legacy
Windows-only COM (``pywin32``) implementation.

The public surface is intentionally identical to the legacy generator so it is
a drop-in replacement::

    PowerPointGenerator().create_presentation(slides, output_path)
"""

from __future__ import annotations

import os

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover - only hit when dependency is missing
    Presentation = None
    Inches = None


class PowerPointGenerator:
    """Render a list of slide dicts into a ``.pptx`` file.

    Each slide is a dict with ``title`` (str) and ``content`` (list). Content
    items are either strings (paragraphs/bullets) or dicts shaped like
    ``{"image": path}`` or ``{"link": href}``.
    """

    def __init__(self, template_path: str | None = None):
        if Presentation is None:
            raise RuntimeError(
                "python-pptx is not installed. Install it with: pip install python-pptx"
            )
        self._template_path = template_path

    def create_presentation(self, slides, output_path):
        """Build the presentation and save it to ``output_path``."""
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        prs = Presentation(self._template_path) if self._template_path else Presentation()

        if not slides:
            print("Warning: No slides to process. The presentation will be empty.")

        for slide in slides or []:
            try:
                self._add_slide(prs, slide)
            except Exception as exc:  # noqa: BLE001 - keep going on a bad slide
                print(f"Error building slide '{slide.get('title', '?')}': {exc}")

        prs.save(output_path)
        print(f"Presentation saved at: {output_path}")
        return output_path

    def _add_slide(self, prs, slide):
        title = slide.get("title", "Untitled Slide")
        contents = slide.get("content", []) or []

        text_items = [c for c in contents if isinstance(c, str) and c.strip()]
        link_items = [c["link"] for c in contents if isinstance(c, dict) and "link" in c]
        image_items = [c["image"] for c in contents if isinstance(c, dict) and "image" in c]

        body_lines = list(text_items) + [f"Link: {href}" for href in link_items]

        # Layout 1 = "Title and Content"; layout 5 = "Title Only".
        layout_index = 1 if body_lines else 5
        layout = prs.slide_layouts[layout_index]
        ppt_slide = prs.slides.add_slide(layout)

        if ppt_slide.shapes.title is not None:
            ppt_slide.shapes.title.text = title

        if body_lines:
            body = self._find_body_placeholder(ppt_slide)
            if body is not None:
                text_frame = body.text_frame
                text_frame.text = body_lines[0]
                for line in body_lines[1:]:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = line

        top = Inches(2.5)
        for image_path in image_items:
            if image_path and os.path.exists(image_path):
                try:
                    ppt_slide.shapes.add_picture(image_path, Inches(1), top, width=Inches(4))
                    top += Inches(3)
                except Exception as exc:  # noqa: BLE001
                    print(f"Error adding image '{image_path}': {exc}")
            elif image_path:
                print(f"Skipping missing image: {image_path}")

    @staticmethod
    def _find_body_placeholder(ppt_slide):
        """Return the content/body placeholder for a slide, if present."""
        for placeholder in ppt_slide.placeholders:
            # idx 0 is the title; the first non-title placeholder is the body.
            if placeholder.placeholder_format.idx != 0:
                return placeholder
        return None
